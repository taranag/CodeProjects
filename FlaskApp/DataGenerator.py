# DataGenerator.py
# This file contains the main function that generates reports
#
# Author: Taran Agnihotri
# Last Updated: 22/6/2022
# Version: 1.0

# Imports
import itertools
import pptx
from pptx.util import Pt
from pptx.util import Inches
from SQLTools import *
import requests
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from PPTXTools import *

# Database Details
host = 'localhost'
user_name = 'taran'
user_password = 'DESKTOP-RMV78RR'
db_name = 'seekapp'

# Constants
myPath = "static/generated/"
logoPath = "SeekLogo.png"
max_table_size = 8
max_table_width = 5
logoLeft, logoTop, logoHeight, logoWidth = 8293608, 18288, 768096, 841248

def fastTranslateToEnglish(session, text):
    '''Use google translate api to automatically translate text to english'''
    '''Requires a session object to be passed in'''
    url = "https://translate.googleapis.com/translate_a/single?client=gtx&sl=auto&tl=en&dt=t&q=" + text
    response = session.get(url)
    response = response.json()
    returnString = ""
    for item in response[0]:
        returnString += item[0]
    return returnString
    
def getCompanies():
    '''Returns a list of all companies in the database'''
    connection = create_server_connection(host, user_name, user_password, db_name)
    query = "SELECT id, name FROM company where status='active'"
    result = execute_query(connection, query)
    return result

def addDownloadData(prs, companyID, groupBy, connection = None):
    '''Adds download data slides to the presentation'''

    # If no database connection is passed in, create a new one
    if connection is None:
        connection = create_server_connection(host, user_name, user_password, db_name)

    # Get all non inactive employees in the company
    # Query for appropriate columns and company ID
    query1 = "select id,{},status from employee where status !='inactive' and companyID={}".format(groupBy, companyID)
    # Results in form (id, groupBy, status)

    # Attempt to execute query and set result1 to the result of the query
    result1 = execute_queryNoTime(connection, query1)

    # create dictionary of employee dictionaries
    employeeDict = {}
    for row in result1:
        employeeDict[row[0]] = {groupBy: row[1], "status": row[2]}
        # format: employeeDict[id] = [dept, status]

    # Test code for employeeDict
    # print("Employee dictionary created with {} employees".format(len(employeeDict)))

    # Create dictionary of unit dictionaries
    activeInactiveByGroup = {}
    for employee in employeeDict.values():
        # If no group is specified, default to "Other"
        if(employee[groupBy] == None):
            employee[groupBy] = "Other"
        currGroup = employee[groupBy].capitalize()
        # If group already exists in dictionary, add employee to respective group
        if currGroup in activeInactiveByGroup:
            if employee["status"] == "active":
                activeInactiveByGroup[currGroup][0] += 1
            else:
                activeInactiveByGroup[currGroup][1] += 1
        # If group does not exist in dictionary, create it
        else:
            if employee["status"] == "active":
                activeInactiveByGroup[currGroup] = [1, 0]
            else:
                activeInactiveByGroup[currGroup] = [0, 1]

    # Test code to print out active/inactive by group
    # print(activeInactiveByGroup)

    # Add download data slides to presentation
    chunkCounter = 0
    while (len(activeInactiveByGroup) > chunkCounter*max_table_size):
        if chunkCounter == 0:
            # Create first slide
            slide = createBlankSlideWithTitle(prs, "Download Status")
        else:
            # Create continuation slide
            slide = createBlankSlideWithTitle(prs, "Download Status Continued")
        
        # Set currChunk to the current chunk of the data dictionary
        currChunk = dict(itertools.islice(activeInactiveByGroup.items(), chunkCounter*max_table_size, chunkCounter*max_table_size + max_table_size))
        
        # Create table
        groupDownloadTable = createTableWithDownloadHeaders(slide, len(currChunk) + 2, 5, Inches(1), Inches(1.5), Inches(8), Inches(2))
        
        # Create lists for efficiency (to avoid iterating over dictionary multiple times)
        currChunkKeys = list(currChunk.keys())
        currChunkValues = list(currChunk.values())

        # Add data to table
        for i in range (0, len(currChunk)):
            # Add row to table
            # Set first column to group name
            groupDownloadTable.rows[i+1].cells[0].text = currChunkKeys[i]
            # Set second column to number of active employees (downloaded)
            groupDownloadTable.rows[i+1].cells[1].text = str(currChunkValues[i][0])
            # Set third column to number of inactive employees (yet to download)
            groupDownloadTable.rows[i+1].cells[2].text = str(currChunkValues[i][1])

            # Set fourth column to total number of employees
            groupDownloadTable.rows[i+1].cells[3].text = str(currChunkValues[i][0] + currChunkValues[i][1])

            # Set fifth column to percentage of active employees
            if currChunkValues[i][0] + currChunkValues[i][1] == 0:
                groupDownloadTable.rows[i+1].cells[4].text = "0%"
            else:
                groupDownloadTable.rows[i+1].cells[4].text = str(round((currChunkValues[i][0])/(currChunkValues[i][0] + currChunkValues[i][1])*100)) + "%"
        
        # Add total row to table
        groupDownloadTable.rows[len(currChunk)+1].cells[0].text = "Total"
        # Set second column to total number of active employees (downloaded)
        groupDownloadTable.rows[len(currChunk)+1].cells[1].text = str(sum(currChunkValues[i][0] for i in range(len(currChunk))))
        # Set third column to total number of inactive employees (yet to download)
        groupDownloadTable.rows[len(currChunk)+1].cells[2].text = str(sum(currChunkValues[i][1] for i in range(len(currChunk))))
        # Set fourth column to total number of employees
        groupDownloadTable.rows[len(currChunk)+1].cells[3].text = str(sum(currChunkValues[i][0] + currChunkValues[i][1] for i in range(len(currChunk))))
        # Set fifth column to total percentage of active employees
        groupDownloadTable.rows[len(currChunk)+1].cells[4].text = str(round((sum(currChunkValues[i][0] for i in range(len(currChunk))))/(sum(currChunkValues[i][0] + currChunkValues[i][1] for i in range(len(currChunk))))*100)) + "%"
        
        # Increment chunk counter
        chunkCounter += 1
        




def createValueTable(slide, answerDictByGroup, optionList):
    '''Creates a table with the values of the answerDictByGroup dictionary'''

    # If answerDictByGroup longer than 4, create multiple tables
    if len(answerDictByGroup) > 4:
        createDoubleValueTable(slide, answerDictByGroup, optionList)
        return

    # Table positional constants
    left = Inches(2.75)
    top = Inches(2.5)
    width = Inches(7)
    height = Inches(2)

    # Create table
    table = slide.shapes.add_table(len(optionList) + 2, len(answerDictByGroup) + 1, left, top, width, height).table

    # Add headers to table
    table.rows[0].cells[0].text = "Response"
    # Create lists for efficiency (to avoid iterating over dictionary multiple times)
    answerDictByGroupKeyList = list(answerDictByGroup.keys())
    answerDictByGroupValueList = list(answerDictByGroup.values())

    # Set row headers to group names
    for i in range(len(answerDictByGroup)):
        table.rows[0].cells[i+1].text = answerDictByGroupKeyList[i]

    # Set column headers to option names and add data to table
    for i in range(len(optionList)):
        # Set column header to option name
        table.rows[i+1].cells[0].text = optionList[i]
        # Iterate through answerDictByGroup and add data to table
        for j in range(len(answerDictByGroup)):
            table.rows[i+1].cells[j+1].text = str(answerDictByGroup[answerDictByGroupKeyList[j]][optionList[i]])
    
    # Add total row to table
    table.rows[len(optionList)+1].cells[0].text = "Total"
    # Iterate through answerDictByGroupValueList and add total data to table
    for j in range(len(answerDictByGroup)):
        table.rows[len(optionList)+1].cells[j+1].text = str(sum(answerDictByGroupValueList[j].values()))

def createPercentValueTable(slide, answerDictByGroup, optionList):
    '''Creates a table with the percent values of the answerDictByGroup dictionary'''

    # If answerDictByGroup longer than 4, create multiple tables
    if len(answerDictByGroup) > 4:
        createDoublePercentValueTable(slide, answerDictByGroup, optionList)
        return

    # Table positional constants
    left = Inches(2.75)
    top = Inches(2.5)
    width = Inches(7)
    height = Inches(2)

    # Create table
    table = slide.shapes.add_table(len(optionList) + 2, len(answerDictByGroup) + 1, left, top, width, height).table

    # Add headers to table
    table.rows[0].cells[0].text = "Response"
    # Create lists for efficiency (to avoid iterating over dictionary multiple times)
    answerDictByGroupKeyList = list(answerDictByGroup.keys())
    answerDictByGroupValueList = list(answerDictByGroup.values())

    # Set row headers to group names
    for i in range(len(answerDictByGroup)):
        table.rows[0].cells[i+1].text = answerDictByGroupKeyList[i]

    # Set column headers to option names and add data to table by percent
    for i in range(len(optionList)):
        # Set column header to option name
        table.rows[i+1].cells[0].text = optionList[i]

        # Iterate through answerDictByGroup and add data to table by percent
        for j in range(len(answerDictByGroup)):
            try:
                table.rows[i+1].cells[j+1].text = str(round((answerDictByGroup[answerDictByGroupKeyList[j]][optionList[i]])/sum(answerDictByGroupValueList[j].values())*100)) + "%"
            except ZeroDivisionError:
                table.rows[i+1].cells[j+1].text = "0%"

    # Add total row to table
    table.rows[len(optionList)+1].cells[0].text = "Total"
    # Iterate through answerDictByGroupValueList and add total data to table by percent
    for j in range(len(answerDictByGroup)):
        try: 
            table.rows[len(optionList)+1].cells[j+1].text = str(round(sum(answerDictByGroupValueList[j].values())/sum(answerDictByGroupValueList[j].values())*100)) + "%"
        except ZeroDivisionError:
            table.rows[len(optionList)+1].cells[j+1].text = "0%"

def createDoublePercentValueTable(slide, answerDictByGroup, optionList):
    # Split dictionary into two dictionaries of equal length
    answerDictByGroup1 = {}
    answerDictByGroup2 = {}
    # Create lists for efficiency (to avoid iterating over dictionary multiple times)
    answerDictByGroupKeyList = list(answerDictByGroup.keys())
    answerDictByGroupValueList = list(answerDictByGroup.values())
    # Split dictionary into two dictionaries
    for i in range(len(answerDictByGroup)):
        if i % 2 == 0:
            answerDictByGroup1[answerDictByGroupKeyList[i]] = answerDictByGroupValueList[i]
        else:
            answerDictByGroup2[answerDictByGroupKeyList[i]] = answerDictByGroupValueList[i]

    # Set current answer dictionary to first dictionary
    answerDictByGroup = answerDictByGroup1

    # First table positional constants
    left = Inches(2.75)
    top = Inches(1.5)
    width = Inches(7)
    height = Inches(2)

    # Create first table
    table = slide.shapes.add_table(len(optionList) + 2, len(answerDictByGroup) + 1, left, top, width, height).table

    # Add headers to table
    table.rows[0].cells[0].text = "Response"

    # Create lists for efficiency (to avoid iterating over dictionary multiple times)
    answerDictByGroupKeyList = list(answerDictByGroup.keys())
    answerDictByGroupValueList = list(answerDictByGroup.values())

    # Set row headers to group names
    for i in range(len(answerDictByGroup)):
        table.rows[0].cells[i+1].text = answerDictByGroupKeyList[i]

    # Set column headers to option names and add data to table by percent
    for i in range(len(optionList)):
        # Set column header to option name
        table.rows[i+1].cells[0].text = optionList[i]
        # Iterate through answerDictByGroup and add data to table by percent
        for j in range(len(answerDictByGroup)):
            try:
                table.rows[i+1].cells[j+1].text = str(round((answerDictByGroup[answerDictByGroupKeyList[j]][optionList[i]])/sum(answerDictByGroupValueList[j].values())*100)) + "%"
            except ZeroDivisionError:
                table.rows[i+1].cells[j+1].text = "0%"
    
    # Add total row to table
    table.rows[len(optionList)+1].cells[0].text = "Total"
    # Iterate through answerDictByGroupValueList and add total data to table by percent
    for j in range(len(answerDictByGroup)):
        try:
            table.rows[len(optionList)+1].cells[j+1].text = str(round(sum(answerDictByGroupValueList[j].values())/sum(answerDictByGroupValueList[j].values())*100)) + "%"
        except ZeroDivisionError:
            table.rows[len(optionList)+1].cells[j+1].text = "0%"

    # Set table font size
    setTableFontSize(table, 16)

    # Do the whole thing over again for the second dictionary
    # Set current answer dictionary to second dictionary
    answerDictByGroup = answerDictByGroup2

    # Second able positional constants
    left = Inches(2.75)
    top = Inches(4.25)
    width = Inches(7)
    height = Inches(2)

    # Create second table
    table2 = slide.shapes.add_table(len(optionList) + 2, len(answerDictByGroup) + 1, left, top, width, height).table

    # Add headers to table
    table2.rows[0].cells[0].text = "Response"

    # Create lists for efficiency (to avoid iterating over dictionary multiple times)
    answerDictByGroupKeyList = list(answerDictByGroup.keys())
    answerDictByGroupValueList = list(answerDictByGroup.values())

    # Set row headers to group names
    for i in range(len(answerDictByGroup)):
        table2.rows[0].cells[i+1].text = answerDictByGroupKeyList[i]

    # Set column headers to option names and add data to table by percent
    for i in range(len(optionList)):
        # Set column header to option name
        table2.rows[i+1].cells[0].text = optionList[i]
        # Iterate through answerDictByGroup and add data to table by percent
        for j in range(len(answerDictByGroup)):
            try:
                table2.rows[i+1].cells[j+1].text = str(round((answerDictByGroup[answerDictByGroupKeyList[j]][optionList[i]])/sum(answerDictByGroupValueList[j].values())*100)) + "%"
            except ZeroDivisionError:
                table2.rows[i+1].cells[j+1].text = "0%"

    # Add total row to table        
    table2.rows[len(optionList)+1].cells[0].text = "Total"
    # Iterate through answerDictByGroupValueList and add total data to table by percent
    for j in range(len(answerDictByGroup)):
        try:
            table2.rows[len(optionList)+1].cells[j+1].text = str(round(sum(answerDictByGroupValueList[j].values())/sum(answerDictByGroupValueList[j].values())*100)) + "%"
        except ZeroDivisionError:
            table2.rows[len(optionList)+1].cells[j+1].text = "0%"
    
    # Set table font size
    setTableFontSize(table2, 16)


def createDoubleValueTable(slide, answerDictByGroup, optionList):
    # Split dictionary into two dictionaries of equal length
    answerDictByGroup1 = {}
    answerDictByGroup2 = {}
    # Create lists for efficiency (to avoid iterating over dictionary multiple times)
    answerDictByGroupKeyList = list(answerDictByGroup.keys())
    answerDictByGroupValueList = list(answerDictByGroup.values())
    # Split dictionary into two dictionaries
    for i in range(len(answerDictByGroup)):
        if i % 2 == 0:
            answerDictByGroup1[answerDictByGroupKeyList[i]] = answerDictByGroupValueList[i]
        else:
            answerDictByGroup2[answerDictByGroupKeyList[i]] = answerDictByGroupValueList[i]
    # Create tables
    # Set current answer dictionary to first dictionary
    answerDictByGroup = answerDictByGroup1

    # First table positional constants
    left = Inches(2.75)
    top = Inches(1.5)
    width = Inches(7)
    height = Inches(2)

    # Create first table
    table = slide.shapes.add_table(len(optionList) + 2, len(answerDictByGroup) + 1, left, top, width, height).table

    # Add headers to table
    table.rows[0].cells[0].text = "Response"

    # Create lists for efficiency (to avoid iterating over dictionary multiple times)
    answerDictByGroupKeyList = list(answerDictByGroup.keys())
    answerDictByGroupValueList = list(answerDictByGroup.values())

    # Set row headers to group names
    for i in range(len(answerDictByGroup)):
        table.rows[0].cells[i+1].text = answerDictByGroupKeyList[i]

    # Set column headers to option names and add data to table
    for i in range(len(optionList)):
        # Set column header to option name
        table.rows[i+1].cells[0].text = optionList[i]
        # Iterate through answerDictByGroup and add data to table
        for j in range(len(answerDictByGroup)):
            table.rows[i+1].cells[j+1].text = str(answerDictByGroup[answerDictByGroupKeyList[j]][optionList[i]])
    
    # Add total row to table
    table.rows[len(optionList)+1].cells[0].text = "Total"
    for j in range(len(answerDictByGroup)):
        table.rows[len(optionList)+1].cells[j+1].text = str(sum(answerDictByGroupValueList[j].values()))

    # Set first table font size
    setTableFontSize(table, 16)

    # Do the whole thing over again for the second dictionary
    # Set current answer dictionary to second dictionary
    answerDictByGroup = answerDictByGroup2

    # Second table positional constants
    left = Inches(2.75)
    top = Inches(4.25)
    width = Inches(7)
    height = Inches(2)

    # Create second table
    table2 = slide.shapes.add_table(len(optionList) + 2, len(answerDictByGroup) + 1, left, top, width, height).table

    # Add headers to table
    table2.rows[0].cells[0].text = "Response"

    # Create lists for efficiency (to avoid iterating over dictionary multiple times)
    answerDictByGroupKeyList = list(answerDictByGroup.keys())
    answerDictByGroupValueList = list(answerDictByGroup.values())

    # Set row headers to group names
    for i in range(len(answerDictByGroup)):
        table2.rows[0].cells[i+1].text = answerDictByGroupKeyList[i]

    # Set column headers to option names and add data to table
    for i in range(len(optionList)):
        # Set column header to option name
        table2.rows[i+1].cells[0].text = optionList[i]
        # Iterate through answerDictByGroup and add data to table
        for j in range(len(answerDictByGroup)):
            table2.rows[i+1].cells[j+1].text = str(answerDictByGroup[answerDictByGroupKeyList[j]][optionList[i]])
    
    # Add total row to table
    table2.rows[len(optionList)+1].cells[0].text = "Total"
    # Iterate through answerDictByGroupValueList and add total data to table
    for j in range(len(answerDictByGroup)):
        table2.rows[len(optionList)+1].cells[j+1].text = str(sum(answerDictByGroupValueList[j].values()))

    # Set second table font size
    setTableFontSize(table2, 16)


def createValuePieChart(slide, answerDictByGroup, optionList):
    '''Create a pie chart for the given answer dictionary and option list.'''
    # Create new chart data object
    chart_data = ChartData()

    # Create chart categories list
    chartCategoriesList = []
    for i in range(len(optionList)):
        # Add option name to chart categories list
        # If option name is longer than 20 characters, truncate it
        if len(optionList[i]) > 20:
            chartCategoriesList.append(optionList[i][:17] + "...")
        else:
            chartCategoriesList.append(optionList[i])
        
    # Set chart categories
    chart_data.categories = chartCategoriesList
    # Create list of answer data keys
    answerDictByGroupKeyList = list(answerDictByGroup.keys())
    # Create dictionary of answer data
    numAnswersDict = {}
    # Iterate through answerDictByGroup and add data to dictionary
    for i in range(len(optionList)):
        numAnswersDict[optionList[i]] = sum(answerDictByGroup[answerDictByGroupKeyList[j]][optionList[i]] for j in range(len(answerDictByGroup)))
    
    # Iterate through numAnswersDict and convert to percentages
    answerPercentageDict = {}
    for i in range(len(optionList)):
        try:
            answerPercentageDict[optionList[i]] = numAnswersDict[optionList[i]] / sum(numAnswersDict.values())
        except ZeroDivisionError:
            answerPercentageDict[optionList[i]] = 0

    # Set chart data
    chart_data.add_series('Response', answerPercentageDict.values())

    # Set height based on amount of options
    if len(optionList) > 4:
        height = Inches(3.5)
    else:
        height = Inches(3)
    
    # Create chart
    pieChart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(0), Inches(2.5), Inches(3), height, chart_data).chart
    # Set chart style
    pieChart.chart_style = 26
    pieChart.has_legend = True
    pieChart.legend.position = pptx.enum.chart.XL_LEGEND_POSITION.BOTTOM
    pieChart.legend.include_in_layout = False
    pieChart.plots[0].has_data_labels = True
    data_labels = pieChart.plots[0].data_labels
    data_labels.number_format = '0%'
    data_labels.position = pptx.enum.chart.XL_LABEL_POSITION.OUTSIDE_END


def addLearnData(prs, companyID, groupBy, startDate, endDate, connection = None):
    '''Add learn data to the given presentation.'''

    # If connection is not given, create a new connection
    if connection is None:
        connection = create_server_connection(host, user_name, user_password, db_name)

    # Get all active employees in the company
    query1 = "select {},count(*) from employee where status ='active' and companyID={} group by {};".format(groupBy, companyID, groupBy)
    # Results in form (id, groupBy, status)
    # Get course data with group and count
    query2 = '''select c.display_name,e.{},count(*) from score s
left join employee e on e.id = s.userid
left join course c on c.course_id = s.course_id
where s.userid in (select id from employee where companyId = {})
and s.date >='{}' and s.date <='{}'
group by c.display_name,e.{};'''.format(groupBy, companyID, str(startDate), str(endDate), groupBy)
    # Results in form (course_name, groupBy, count)

    # Attempt to execute queries and get results
    result1 = execute_queryNoTime(connection, query1)
    result2 = execute_queryNoTime(connection, query2)

    totalEmployeesByGroup = {}
    for group in result1:
        try:
            currGroup = group[0].capitalize()
        except:
            currGroup = "Other"
        if currGroup in totalEmployeesByGroup:
            totalEmployeesByGroup[currGroup] += group[1]
        else:
            totalEmployeesByGroup[currGroup] = group[1]

    # Test code for totalEmployeesByGroup   
    # print("Total employees by group: {}".format(totalEmployeesByGroup))
    
    dataDict = {}
    courseList = []
    for row in result2:
        try:
            currGroup = row[1].capitalize()
        except:
            currGroup = "Other"
        if row[0] not in courseList:
            courseList.append(row[0])
        if currGroup in dataDict:
            dataDict[currGroup][row[0]] = row[2]
        else:
            dataDict[currGroup] = {row[0]: row[2]}

    for group in totalEmployeesByGroup:
        if group not in dataDict:
            dataDict[group] = {}
    totalEmployeesLearningByGroup = {}
    #totalEmployeesLearningByGroup["total"] = 0

    chunkCounter = 0
    while (len(dataDict) > chunkCounter*max_table_size):
        sliceCounter = 0
        while(len(courseList) > sliceCounter*max_table_width):
            if chunkCounter == 0 and sliceCounter == 0:
                slide = createBlankSlideWithTitle(prs, "Learn Status")
            else:
                slide = createBlankSlideWithTitle(prs, "Learn Status Continued")
            # currCourseList = courseList[chunkCounter*max_table_size:(chunkCounter+1)*max_table_size]
            currCourseList = courseList[sliceCounter*max_table_width:(sliceCounter+1)*max_table_width]
            currChunk = dict(itertools.islice(dataDict.items(), chunkCounter*max_table_size, chunkCounter*max_table_size + max_table_size))
            if len(currCourseList) < (max_table_width):
                groupLearnTable = createTableWithFinalLearnHeaders(slide, currCourseList, len(currChunk) + 2, len(currCourseList) + 3)
            else:
                groupLearnTable = createTableWithLearnHeaders(slide, currCourseList, len(currChunk) + 2, len(currCourseList) + 1)
            currChunkKeys = list(currChunk.keys())
            currChunkValues = list(currChunk.values())
            for i in range (0, len(currChunkKeys)):
                totalLearning = 0
                groupLearnTable.rows[i+1].cells[0].text = currChunkKeys[i]
                for j in range (0, len(currCourseList)):
                    if currCourseList[j] in currChunkValues[i]:
                        groupLearnTable.rows[i+1].cells[j+1].text = str(currChunkValues[i][currCourseList[j]])
                        totalLearning += currChunk[currChunkKeys[i]][currCourseList[j]]
                    else:
                        groupLearnTable.rows[i+1].cells[j+1].text = "0"
                if currChunkKeys[i] in totalEmployeesLearningByGroup:
                    totalEmployeesLearningByGroup[currChunkKeys[i]] += totalLearning
                else:
                    totalEmployeesLearningByGroup[currChunkKeys[i]] = totalLearning
                #totalEmployeesLearningByGroup["total"] += totalLearning
                if len(currCourseList) < (max_table_width):

                    groupLearnTable.rows[i+1].cells[len(currCourseList)+1].text = str(totalEmployeesByGroup[currChunkKeys[i]] - totalEmployeesLearningByGroup[currChunkKeys[i]])
                    groupLearnTable.rows[i+1].cells[len(currCourseList)+2].text = str(totalEmployeesByGroup[currChunkKeys[i]])
            groupLearnTable.rows[len(currChunkKeys)+1].cells[0].text = "Total"

            companyLearning = 0
            for i in range (0, len(currCourseList)):
                totalLearning = 0
                for j in range (0, len(currChunkKeys)):
                    if currCourseList[i] in currChunkValues[j]:
                        totalLearning += currChunk[currChunkKeys[j]][currCourseList[i]]
                companyLearning += totalLearning
                groupLearnTable.rows[len(currChunkKeys)+1].cells[i+1].text = str(totalLearning)
            
            # Runs on last slide of slice
            if len(currCourseList) < (max_table_width):
                currLearning = 0
                for i in range (0, len(currChunkKeys)):
                    currLearning += totalEmployeesLearningByGroup[currChunkKeys[i]]
                currEmployeeCount = 0
                for i in range (0, len(currChunkKeys)):
                    currEmployeeCount += totalEmployeesByGroup[currChunkKeys[i]]
                groupLearnTable.rows[len(currChunkKeys)+1].cells[len(currCourseList)+1].text = str(currEmployeeCount - currLearning)
                groupLearnTable.rows[len(currChunkKeys)+1].cells[len(currCourseList)+2].text = str(currEmployeeCount)
            sliceCounter += 1

        chunkCounter += 1

def addValueData(prs, companyID, groupBy, startDate, endDate, percentage=1, connection = None):
    '''Add value data to the given presentation.'''

    # If no connection is given, create a new one
    if connection is None:
        connection = create_server_connection(host, user_name, user_password, db_name)

    # Get all questions asked for a company during a given time period
    query1 = "select id,question,option1,option2,option3,option4,option5,employees_type_data,actual_schedule_time from surveyquestions where compid = {} and actual_schedule_time >= '{}' and actual_schedule_time <= '{}' and scheduled_count<>0;".format(companyID, startDate, endDate)
    result1 = execute_queryNoTime(connection, query1)

    query3 = "select {},count(*) from employee where companyId={} and status='active' group by {};".format(groupBy, companyID, groupBy)
    result3 = execute_queryNoTime(connection, query3)
    employeeGroupList = []
    for row3 in result3:
        group = row3[0]
        if group == None or group == "":
            group = "Other"
        else:
            group = group.capitalize()
        employeeGroupList.append(group)

    # Remove duplicates from employeeGroupList
    employeeGroupList = list(set(employeeGroupList))

    allDataByQuestion = []

    s = requests.Session()
    for row1 in result1:
        questionID = row1[0]
        question = row1[1]
        if not (row1[7] == "English" or row1[7] == "FC" or row1[7] == "DP-English"):
            if question is not None:
                # question = translateToEnglishDos(question)
                question = fastTranslateToEnglish(s, question)
        # check if question is null
        if question is None:
            tempQuery = "select content_file from surveyquestions where id = {}".format(questionID)
            tempResult = execute_queryNoTime(connection, tempQuery)
            question = tempResult[0][0]
            #question = "https://backoffice.seek-app.com//storage/" + question
            #question = processURL(question).replace('\\n', " ")
            
        query2 = "select e.{},n.answer from notifications n left join employee e on n.empId=e.id where message_id={} and answer is not null;".format(groupBy, questionID)
        result2 = execute_queryNoTime(connection, query2)

        optionList = []
        for i in range(2,7):
            if row1[i] != None and row1[i] != "":
                optionList.append(row1[i])

        answerDictByGroup = {}
        for group in employeeGroupList:
            answerDictByGroup[group] = {}
            for option in optionList:
                answerDictByGroup[group][option] = 0
        
        
        for row2 in result2:
            answerDictByGroup[row2[0].capitalize()][row2[1]] += 1

        allDataByQuestion.append([question, answerDictByGroup, optionList, row1[7], row1[8]])
            
    # Remove duplicates
    currLength = len(allDataByQuestion) - 1
    idx = 0
    while idx < currLength:
        try:
            isSame = allDataByQuestion[idx][4] == allDataByQuestion[idx+1][4] 
            # and allDataByQuestion[idx][3] != allDataByQuestion[idx+1][3]
            if isSame:
                firstAnswerDict = allDataByQuestion[idx][1]
                secondAnswerDict = allDataByQuestion[idx+1][1]
                firstAnswerDictKeyList = list(firstAnswerDict.keys())
                for i in range(len(firstAnswerDictKeyList)):
                    currGroup = firstAnswerDictKeyList[i]
                    k = 0
                    secondAnswerCurrGroupValueList = list(secondAnswerDict[currGroup].values())
                    for j in firstAnswerDict[currGroup].keys():
                        firstAnswerDict[currGroup][j] += secondAnswerCurrGroupValueList[k]
                        k += 1
                allDataByQuestion.pop(idx+1)
                currLength -= 1
                continue
        except:
            print("False duplicate found at index: " + str(idx))
            print("Confused questions: " + str(allDataByQuestion[idx][0]) + " and " + str(allDataByQuestion[idx+1][0]))
            pass
        idx += 1

    # Add Slides
    for i in range(len(allDataByQuestion)):
        currQuestion = allDataByQuestion[i]
        question = currQuestion[0]
        answerDictByGroup = currQuestion[1]
        optionList = currQuestion[2]
        # Create a new slide for each question
        # print(questionID, question)
        if len(question) > 180:
            slide = createBlankSlideWithTitle(prs, question, 18)
        else:
            slide = createBlankSlideWithTitle(prs, question, 25)

        # Create a table for the question
        if percentage == 2:
            createPercentValueTable(slide, answerDictByGroup, optionList)
        else:
            createValueTable(slide, answerDictByGroup, optionList)
        # Create a pie chart for the question
        createValuePieChart(slide, answerDictByGroup, optionList)
    

def generateFullReport(companyID, filename, groupBy, startDate, endDate, options):
    startTime = time.time()
    prs = pptx.Presentation()
    connection = create_server_connection(host, user_name, user_password, db_name)
    if (options[0] == 1):
        if sum(options) == 4:
            titleSlide = createTitleSlide(prs, "Full Report for " + companyID, startDate, endDate)
        else:
            titleText = "Report for " + companyID + " with "
            if options[1] == 1:
                titleText += "download data, "
            if options[2] == 1:
                titleText += "learn data, "
            if options[3] == 1 or options[3] == 2:
                titleText += "value data, "
            titleText = titleText[:-2]
            titleSlide = createTitleSlide (prs, titleText, startDate, endDate)
    if (options[1] != 0):
        addDownloadData(prs, companyID, groupBy, connection)
        print("Download data added after {} seconds".format(time.time() - startTime))
    if (options[2] != 0):
        addLearnData(prs, companyID, groupBy, startDate, endDate, connection)
        print("Learn data added after {} seconds".format(time.time() - startTime))
    if (options[3] != 0):
        addValueData(prs, companyID, groupBy, startDate, endDate, options[3], connection)
        print("Value data added after {} seconds".format(time.time() - startTime))


    print("PPTX file creation took {} seconds".format(time.time() - startTime))
    saved = False
    try:
        prs.save(myPath + filename + ".pptx")
        saved = True
        print("File saved as " + myPath + filename + ".pptx")
        return (myPath + filename + ".pptx")
    except:
        number = filename[-1]
        try:
            number = int(number)
        except:
            number = 1
            filename = filename + str(number)
    while(saved == False):
        try:
            prs.save(myPath + filename[:-1] + str(number) + ".pptx")
            print("File saved as " + myPath + filename[:-1] + str(number) + ".pptx")
            saved = True
        except Exception as e:
            print("Save failed. Error: {} Trying again.".format(e))
            number += 1
        if(number == 10):
            print("File could not be saved.")
            return None
    return (myPath + filename[:-1] + str(number) + ".pptx")

generateFullReport("51", "51Test1", "dept", "2022-06-01", "2022-06-14", (1, 1, 1, 2))