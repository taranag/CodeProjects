import pptx
from pptx.util import Pt
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


logoLeft, logoTop, logoHeight, logoWidth = 8293608, 18288, 768096, 841248
logoPath = "SeekLogo.png"
letterPath = "letter.png"

def SubElement(parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

def setShapeTransparency(shape, alpha):
    """ Set the transparency (alpha) of a shape"""
    ts = shape.fill._xPr.solidFill
    sF = ts.get_or_change_to_srgbClr()
    sE = SubElement(sF, 'a:alpha', val=str(alpha))

def setTableFontSize(table, fontSize):
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(fontSize)

def createTitleSlide(prs, title, startDate, endDate):
    """
    Creates a blank slide with a title and logo image
    """
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    title = title.replace("\n", " ")

    titleShape = slide.shapes.title
    titleShape.text = title

    subtitleShape = slide.shapes[1]
    subtitleShape.text = "From " + startDate + " to " + endDate

    titleTextFrame = slide.shapes[0].text_frame
    titleTextFrame.paragraphs[0].runs[0].font.bold = True

    logo = slide.shapes.add_picture(logoPath, pptx.util.Inches(8.45), logoTop, pptx.util.Inches(1.54), pptx.util.Inches(1.42))
    return slide

def createEndSlide(prs, text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(56, 87, 36)
    letter = slide.shapes.add_picture(letterPath, pptx.util.Inches(1.95), pptx.util.Inches(2.71), pptx.util.Inches(1.94), pptx.util.Inches(1.47))
    left = Inches(1.9)
    top = Inches(1.5)
    width = Inches(2)
    height = Inches(0.7)
    txBox1 = slide.shapes.add_textbox(left, top, width, height)
    thankBox = txBox1.text_frame
    thankBox.text = "Thank You"
    thankBox.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
    thankBox.paragraphs[0].runs[0].font.size = Pt(32)
    
    left = Inches(5)
    top = Inches(4)
    width = Inches(3)
    height = Inches(0.7)
    txBox2 = slide.shapes.add_textbox(left, top, width, height)
    linkBox = txBox2.text_frame
    linkBox.text = text
    linkBox.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
    linkBox.paragraphs[0].runs[0].font.size = Pt(32)

    left = Inches(4.83)
    top = Inches(2.15)
    width = Inches(4.03)
    height = Inches(4.4)
    squareAround = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    squareAround.fill.solid()
    squareAround.fill.fore_color.rgb = RGBColor(255,255,255)
    
    setShapeTransparency(squareAround,0)


def createBlankSlideWithTitle(prs, title, fontSize=44, hLink=False, hLinkText=""):
    """
    Creates a blank slide with a title and logo image
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    title = title.replace("\r\n", " ")

    titleShape = slide.shapes.title
    titleShape.text = title

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame

    text_frame.paragraphs[0].runs[0].font.size = Pt(fontSize)
    text_frame.paragraphs[0].runs[0].font.bold = True
    if hLink:
        text_frame.paragraphs[0].runs[0].hyperlink.address = hLinkText

    left = top = width = height = pptx.util.Inches(1)
    logo = slide.shapes.add_picture(logoPath, logoLeft, logoTop, logoWidth, logoHeight)
    slide.shapes._spTree.remove(logo._element)
    slide.shapes._spTree.insert(2, logo._element)
    return slide

def createTableWithDownloadHeaders(slide, rows, cols, left, top, width, height):
    """
    Creates a table on a slide.
    """
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    table.columns[2].width = Inches(2)
    table.columns[0].width = Inches(2)
    table.columns[3].width = Inches(1)
    table.rows[0].cells[1].text = "Downloaded"
    table.rows[0].cells[2].text = "Yet to Download"
    table.rows[0].cells[3].text = "Total"
    table.rows[0].cells[4].text = "% Download"
    return table

def createTableWithLearnHeaders(slide, headers, rows, cols):
    """
    Creates a table on a slide.
    """
    left = Inches(.2)
    top = Inches(1.5)
    width = Inches(9.6)
    height = Inches(2)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    # table.table.columns[2].width = Inches(2)
    # table.table.columns[0].width = Inches(2)
    #table.table.columns[cols-1].width = Inches(1)
    for i in range(len(headers)):
        table.rows[0].cells[i+1].text = headers[i]
    return table

def createTableWithFinalLearnHeaders(slide, headers, rows, cols):
    """
    Creates a table on a slide.
    """
    left = Inches(.2)
    top = Inches(1.5)
    width = Inches(9.6)
    height = Inches(2)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    # table.table.columns[2].width = Inches(2)
    # table.table.columns[0].width = Inches(2)
    #table.table.columns[cols-1].width = Inches(1)
    for i in range(len(headers)):
        table.rows[0].cells[i+1].text = headers[i]
    table.rows[0].cells[len(headers)+1].text = "Not on Learn"
    table.rows[0].cells[len(headers)+2].text = "Total"
    # table.columns[len(headers)+2].width = Inches(1)
    return table