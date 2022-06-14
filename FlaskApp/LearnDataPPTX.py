import itertools
import mysql.connector
from mysql.connector import Error
import pptx
from pptx.util import Pt
from pptx.util import Inches
from SQLTools import *

host = 'localhost'
user_name = 'taran'
user_password = 'DESKTOP-RMV78RR'
db_name = 'seekapp'
max_table_size = 8
max_table_width = 5

myPath = "generated/"

logoPath = "SeekLogo.png"

logoLeft = pptx.util.Inches(9.07)
logoTop = pptx.util.Inches(0.02)
logoHeight = pptx.util.Inches(0.84)
logoWidth = pptx.util.Inches(0.92)


def createBlankSlideWithTitle(prs, title):
    """
    Creates a blank slide with a title and logo image
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    titleShape = slide.shapes.title
    titleShape.text = title

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame

    text_frame.paragraphs[0].runs[0].font.size = Pt(44)
    text_frame.paragraphs[0].runs[0].font.bold = True

    left = top = width = height = pptx.util.Inches(1)
    pic = slide.shapes.add_picture(logoPath, logoLeft, logoTop, logoWidth, logoHeight)
    return slide

def createTable(slide, rows, cols, left, top, width, height):
    """
    Creates a table on a slide.
    """
    table = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table.table

def createTableWithLearnHeaders(slide, headers, rows, cols):
    """
    Creates a table on a slide.
    """
    left = Inches(.333)
    top = Inches(1.5)
    width = Inches(9.5)
    height = Inches(2)
    table = slide.shapes.add_table(rows, cols, left, top, width, height)
    # table.table.columns[2].width = Inches(2)
    # table.table.columns[0].width = Inches(2)
    #table.table.columns[cols-1].width = Inches(1)
    for i in range(len(headers)):
        table.table.rows[0].cells[i+1].text = headers[i]
    return table.table

def createTableWithFinalLearnHeaders(slide, headers, rows, cols):
    """
    Creates a table on a slide.
    """
    left = Inches(.333)
    top = Inches(1.5)
    width = Inches(9.5)
    height = Inches(2)
    table = slide.shapes.add_table(rows, cols, left, top, width, height)
    # table.table.columns[2].width = Inches(2)
    # table.table.columns[0].width = Inches(2)
    #table.table.columns[cols-1].width = Inches(1)
    for i in range(len(headers)):
        table.table.rows[0].cells[i+1].text = headers[i]
    table.table.rows[0].cells[len(headers)+1].text = "Not on Learn"
    table.table.rows[0].cells[len(headers)+2].text = "Total"
    return table.table

def generatePPTXLearnData(companyID, filename, groupBy, startDate, endDate):
    startTime = time.time()
    # Connect to MySQL database
    connection = create_server_connection(host, user_name, user_password, db_name)

    # Get all non inactive employees in the company
    # Query for appropriate columns and company ID
    print(startDate)
    print(endDate)
    query1 = "select id,{},status from employee where status !='inactive' and companyID={}".format(groupBy, companyID)
    query2 = '''select c.display_name,e.{},count(*) from score s
left join employee e on e.id = s.userid
left join course c on c.course_id = s.course_id
where s.userid in (select id from employee where companyId = {})
and s.date >='{}' and s.date <='{}'
group by c.display_name,e.{};'''.format(groupBy, companyID, str(startDate), str(endDate), groupBy)

    # Attempt to execute_query and set result1 to the result of the query
    result1 = execute_query(connection, query1)
    result2 = execute_query(connection, query2)
    # create dictionary of employee dictionaries
    employeeDict = {}
    for row in result1:
        employeeDict[row[0]] = {groupBy: row[1], "status": row[2]}
        # format: employeeDict[id] = [dept, status]
    print("Employee dictionary created with {} employees".format(len(employeeDict)))


    totalEmployeesByGroup = {}
    for employee in employeeDict.values():
        if(employee[groupBy] == None):
            employee[groupBy] = "Other"
        currGroup = employee[groupBy].capitalize()
        if currGroup in totalEmployeesByGroup:
            totalEmployeesByGroup[currGroup] += 1
        else:
            totalEmployeesByGroup[currGroup] = 1
    
    dataDict = {}
    courseList = []
    for row in result2:
        if row[0] not in courseList:
            courseList.append(row[0])
        if row[1] in dataDict:
            dataDict[row[1].capitalize()][row[0]] = row[2]
        else:
            dataDict[row[1].capitalize()] = {row[0]: row[2]}

    for group in totalEmployeesByGroup:
        if group not in dataDict:
            dataDict[group] = {}


    prs = pptx.Presentation()
    chunkCounter = 0
    while (len(dataDict) > chunkCounter*max_table_size):
        sliceCounter = 0
        while(len(courseList) > sliceCounter*max_table_width):
            if chunkCounter == 0 and sliceCounter == 0:
                slide = createBlankSlideWithTitle(prs, "Learn Status")
            else:
                slide = createBlankSlideWithTitle(prs, "Learn Status Continued")
            # currCourseList = courseList[chunkCounter*max_table_size:(chunkCounter+1)*max_table_size]
            currCourseList = courseList[sliceCounter*max_table_width:(sliceCounter+1)*max_table_width]
            currChunk = dict(itertools.islice(dataDict.items(), chunkCounter*max_table_size, chunkCounter*max_table_size + max_table_size))
            if len(currCourseList) < (max_table_width - 1):
                groupLearnTable = createTableWithFinalLearnHeaders(slide, currCourseList, len(currChunk) + 2, len(currCourseList) + 3)
            else:
                groupLearnTable = createTableWithLearnHeaders(slide, currCourseList, len(currChunk) + 2, len(currCourseList) + 1)
            currChunkKeys = list(currChunk.keys())
            currChunkValues = list(currChunk.values())
            for i in range (0, len(currChunkKeys)):
                totalLearning = 0
                groupLearnTable.rows[i+1].cells[0].text = currChunkKeys[i]
                for j in range (0, len(currCourseList)):
                    if currCourseList[j] in currChunkValues[i]:
                        groupLearnTable.rows[i+1].cells[j+1].text = str(currChunkValues[i][currCourseList[j]])
                        totalLearning += currChunk[currChunkKeys[i]][currCourseList[j]]
                    else:
                        groupLearnTable.rows[i+1].cells[j+1].text = "0"
                if len(currCourseList) < (max_table_width - 1):
                    groupLearnTable.rows[i+1].cells[len(currCourseList)+1].text = str(totalEmployeesByGroup[currChunkKeys[i]] - totalLearning)
                    groupLearnTable.rows[i+1].cells[len(currCourseList)+2].text = str(totalEmployeesByGroup[currChunkKeys[i]])
            groupLearnTable.rows[len(currChunkKeys)+1].cells[0].text = "Total"

            companyLearning = 0
            for i in range (0, len(currCourseList)):
                totalLearning = 0
                for j in range (0, len(currChunkKeys)):
                    if currCourseList[i] in currChunkValues[j]:
                        totalLearning += currChunk[currChunkKeys[j]][currCourseList[i]]
                companyLearning += totalLearning
                groupLearnTable.rows[len(currChunkKeys)+1].cells[i+1].text = str(totalLearning)

            if len(currCourseList) < (max_table_width - 1):
                groupLearnTable.rows[len(currChunkKeys)+1].cells[len(currCourseList)+1].text = str(sum(totalEmployeesByGroup.values()) - companyLearning)
                groupLearnTable.rows[len(currChunkKeys)+1].cells[len(currCourseList)+2].text = str(sum(totalEmployeesByGroup.values()))
            sliceCounter += 1

        chunkCounter += 1


    saved = False
    try:
        prs.save(myPath + filename + ".pptx")
        saved = True
        print("File saved as " + myPath + filename + ".pptx")
        return (myPath + filename + ".pptx")
    except:
        number = filename[-1]
        try:
            number = int(number)
        except:
            number = 1
            filename = filename + str(number)
    while(saved == False):
        try:
            prs.save(myPath + filename[:-1] + str(number) + ".pptx")
            print("File saved as " + myPath + filename[:-1] + str(number) + ".pptx")
            saved = True
        except Exception as e:
            print("Save failed. Error: {} Trying again.".format(e))
            number += 1
        if(number == 10):
            print("File could not be saved.")
            return None
            break
    print("PPTX file creation took {} seconds".format(time.time() - startTime))
    return (myPath + filename[:-1] + str(number) + ".pptx")


generatePPTXLearnData(51, "Learn5", "dept", "2022-06-01", "2022-06-14")