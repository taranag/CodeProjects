import itertools
import mysql.connector
from mysql.connector import Error
import pptx
from pptx.util import Pt
from pptx.util import Inches
from SQLTools import *

host = 'localhost'
user_name = 'taran'
user_password = 'DESKTOP-RMV78RR'
db_name = 'seekapp'
max_table_size = 8

#myPath = "/Users/Taran Agnihotri/Desktop/CodeProjects/FlaskApp/"
myPath = "generated/"

logoPath = "/Users/Taran Agnihotri/Desktop/CodeProjects/FlaskApp/SeekLogo.png"

logoPath2 = "SeekLogo.png"

logoLeft = pptx.util.Inches(9.07)
logoTop = pptx.util.Inches(0.02)
logoHeight = pptx.util.Inches(0.84)
logoWidth = pptx.util.Inches(0.92)


def createBlankSlideWithTitle(prs, title):
    """
    Creates a blank slide with a title and logo image
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    titleShape = slide.shapes.title
    titleShape.text = title

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame

    text_frame.paragraphs[0].runs[0].font.size = Pt(44)
    text_frame.paragraphs[0].runs[0].font.bold = True

    left = top = width = height = pptx.util.Inches(1)
    pic = slide.shapes.add_picture(logoPath2, logoLeft, logoTop, logoWidth, logoHeight)
    return slide

def createTable(slide, rows, cols, left, top, width, height):
    """
    Creates a table on a slide.
    """
    table = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table.table

def createTableWithDownloadHeaders(slide, rows, cols, left, top, width, height):
    """
    Creates a table on a slide.
    """
    table = slide.shapes.add_table(rows, cols, left, top, width, height)
    table.table.columns[2].width = Inches(2)
    table.table.columns[0].width = Inches(2)
    table.table.columns[3].width = Inches(1)
    table.table.rows[0].cells[1].text = "Downloaded"
    table.table.rows[0].cells[2].text = "Yet to Download"
    table.table.rows[0].cells[3].text = "Total"
    table.table.rows[0].cells[4].text = "% Download"
    return table.table

def generatePPTXDownloadData(companyID, filename, groupBy):
    startTime = time.time()
    # Connect to MySQL database
    connection = create_server_connection(host, user_name, user_password, db_name)

    # Get all non inactive employees in the company
    # Query for appropriate columns and company ID
    query1 = "select id,{},status from employee where status !='inactive' and companyID={}".format(groupBy, companyID)
    # Attempt to execute_query and set result1 to the result of the query
    result1 = execute_query(connection, query1)
    # create dictionary of employee dictionaries
    employeeDict = {}
    for row in result1:
        employeeDict[row[0]] = {groupBy: row[1], "status": row[2]}
        # format: employeeDict[id] = [dept, status]
    print("Employee dictionary created with {} employees".format(len(employeeDict)))

    # Create dictionary of unit dictionaries
    activeInactiveByGroup = {}
    for employee in employeeDict.values():
        if(employee[groupBy] == None):
            employee[groupBy] = "Other"
        currGroup = employee[groupBy].capitalize()
        if currGroup in activeInactiveByGroup:
            if employee["status"] == "active":
                activeInactiveByGroup[currGroup][0] += 1
            else:
                activeInactiveByGroup[currGroup][1] += 1
        else:
            if employee["status"] == "active":
                activeInactiveByGroup[currGroup] = [1, 0]
            else:
                activeInactiveByGroup[currGroup] = [0, 1]

    #print("Active/Inactive dictionary created with {} units".format(len(activeInactiveByGroup)))
    #print(activeInactiveByGroup)

    prs = pptx.Presentation()
    chunkCounter = 0
    while (len(activeInactiveByGroup) > chunkCounter*max_table_size):
        if chunkCounter == 0:
            slide = createBlankSlideWithTitle(prs, "Download Status")
        else:
            slide = createBlankSlideWithTitle(prs, "Download Status Continued")
        currChunk = dict(itertools.islice(activeInactiveByGroup.items(), chunkCounter*max_table_size, chunkCounter*max_table_size + max_table_size))
        groupDownloadTable = createTableWithDownloadHeaders(slide, len(currChunk) + 2, 5, Inches(1), Inches(1.5), Inches(8), Inches(2))
        currChunkKeys = list(currChunk.keys())
        currChunkValues = list(currChunk.values())
        for i in range (0, len(currChunk)):
            groupDownloadTable.rows[i+1].cells[0].text = currChunkKeys[i]
            groupDownloadTable.rows[i+1].cells[1].text = str(currChunkValues[i][0])
            groupDownloadTable.rows[i+1].cells[2].text = str(currChunkValues[i][1])
            groupDownloadTable.rows[i+1].cells[3].text = str(currChunkValues[i][0] + currChunkValues[i][1])
            if currChunkValues[i][0] + currChunkValues[i][1] == 0:
                groupDownloadTable.rows[i+1].cells[4].text = "0%"
            else:
                groupDownloadTable.rows[i+1].cells[4].text = str(round((currChunkValues[i][0])/(currChunkValues[i][0] + currChunkValues[i][1])*100)) + "%"
        groupDownloadTable.rows[len(currChunk)+1].cells[0].text = "Total"
        groupDownloadTable.rows[len(currChunk)+1].cells[1].text = str(sum(currChunkValues[i][0] for i in range(len(currChunk))))
        groupDownloadTable.rows[len(currChunk)+1].cells[2].text = str(sum(currChunkValues[i][1] for i in range(len(currChunk))))
        groupDownloadTable.rows[len(currChunk)+1].cells[3].text = str(sum(currChunkValues[i][0] + currChunkValues[i][1] for i in range(len(currChunk))))
        groupDownloadTable.rows[len(currChunk)+1].cells[4].text = str(round((sum(currChunkValues[i][0] for i in range(len(currChunk))))/(sum(currChunkValues[i][0] + currChunkValues[i][1] for i in range(len(currChunk))))*100)) + "%"
        
        chunkCounter += 1

    saved = False
    try:
        prs.save(myPath + filename + ".pptx")
        saved = True
        print("File saved as " + myPath + filename + ".pptx")
        return (myPath + filename + ".pptx")
    except:
        number = filename[-1]
        try:
            number = int(number)
        except:
            number = 1
            filename = filename + str(number)
    while(saved == False):
        try:
            prs.save(myPath + filename[:-1] + str(number) + ".pptx")
            print("File saved as " + myPath + filename[:-1] + str(number) + ".pptx")
            saved = True
        except Exception as e:
            print("Save failed. Error: {} Trying again.".format(e))
            number += 1
        if(number == 10):
            print("File could not be saved.")
            return None
            break
    print("PPTX file creation took {} seconds".format(time.time() - startTime))
    return (myPath + filename[:-1] + str(number) + ".pptx")


#generatePPTXDownloadData(92, "Test1", "level")